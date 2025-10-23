"""
═══════════════════════════════════════════════════════════════════════════════
GJENERUES AUTOMATIK I PREZANTIMIT POWERPOINT
═══════════════════════════════════════════════════════════════════════════════

Krijon prezantim profesional PowerPoint me:
- Titull dhe hyrje
- 12+ figura te gjeneruara
- Tabela me rezultate
- Statistika kryesore
- Konkluzioni

Perdor python-pptx library per automatizim total
═══════════════════════════════════════════════════════════════════════════════
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime

def create_title_slide(prs):
    """Krijon slide-in e titullit"""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ANALIZA E AVANCUAR FINANCIARE"
    subtitle.text = ("Multi-Asset Portfolio Analysis\n"
                     "Algoritme të Sofistikuara PySpark\n\n"
                     "Universiteti i Prishtinës\n"
                     "Republika e Kosovës\n\n"
                     f"{datetime.now().strftime('%B %Y')}")
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return prs

def create_agenda_slide(prs):
    """Krijon slide-in e agjendës"""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "AGJENDA E PREZANTIMIT"
    
    tf = body.text_frame
    tf.text = "Hyrje dhe Objektivat"
    
    topics = [
        "Dataset-i Multi-Asset (12 assets)",
        "Algoritmet e Implementuara (10)",
        "Time Series Analysis",
        "Volatility dhe Risk Management",
        "Portfolio Performance",
        "Correlation Analysis",
        "Statistical Distributions",
        "Candlestick Charts",
        "3D Visualizations",
        "Heatmaps dhe Matricat",
        "Konkluzioni dhe Rekommandimet"
    ]
    
    for topic in topics:
        p = tf.add_paragraph()
        p.text = topic
        p.level = 1
    
    return prs

def add_algorithms_slide(prs):
    """Slide me algoritmet e implementuara"""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "ALGORITMET E SOFISTIKUARA - PYSPARK"
    
    tf = body.text_frame
    tf.text = "1. PERCENTILE_APPROX - Kuantilet (P25, P50, P75, P95, P99)"
    
    algorithms = [
        "2. STDDEV_SAMP - Sample Standard Deviation",
        "3. LAG - Time Series Operations (1, 5, 10, 20 ditë)",
        "4. MOVING AVERAGE - MA 7d, 30d, 90d",
        "5. GROUPBY AGGREGATION - MapReduce Style",
        "6. WINDOW FUNCTIONS - rank, dense_rank, percent_rank, ntile",
        "7. NULL COUNTING - Data Completeness Analysis",
        "8. DATA PROFILING - Mean, Variance, Skewness, Kurtosis",
        "9. RETURNS CALCULATION - Daily, Weekly, Monthly",
        "10. VOLATILITY COMPUTATION - Historical & Rolling"
    ]
    
    for algo in algorithms:
        p = tf.add_paragraph()
        p.text = algo
        p.level = 0
        p.font.size = Pt(16)
    
    return prs

def add_image_slide(prs, title_text, image_path):
    """Shton slide me figure"""
    blank_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title_text
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add image if exists
    if os.path.exists(image_path):
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        
        try:
            slide.shapes.add_picture(image_path, left, top, width=width)
        except Exception as e:
            # Add error text if image fails
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            txBox2.text_frame.text = f"Figura nuk mund të ngarkohet: {os.path.basename(image_path)}"
    else:
        # Add placeholder text
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox2.text_frame.text = f"Figura: {os.path.basename(image_path)}\n(Do të gjenerohet nga analiza)"
    
    return prs

def create_conclusion_slide(prs):
    """Slide i konkluzioni"""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "KONKLUZIONI DHE ARRITJET"
    
    tf = body.text_frame
    tf.text = "✓ Dataset Multi-Asset me 12 instrumente financiare"
    
    conclusions = [
        "✓ 10 Algoritme të sofistikuara PySpark të implementuara",
        "✓ 12+ Vizualizime profesionale të gjeneruara",
        "✓ Analiza e plotë e Time Series",
        "✓ Risk Management dhe Volatility Analysis",
        "✓ Portfolio Optimization",
        "✓ Correlation dhe Diversification Analysis",
        "✓ Statistical Profiling i detajuar",
        "✓ Automated PowerPoint Generation",
        "✓ Nivel Doktoral - Gatshme për publikim"
    ]
    
    for conclusion in conclusions:
        p = tf.add_paragraph()
        p.text = conclusion
        p.level = 0
        p.font.size = Pt(18)
    
    return prs

def main():
    print("="*80)
    print("GJENERUES AUTOMATIK I POWERPOINT PREZANTIMIT")
    print("="*80)
    print()
    
    # Krijo prezantimin
    print("[1/4] Duke krijuar prezantimin PowerPoint...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    print("[2/4] Duke shtuar slides...")
    prs = create_title_slide(prs)
    print("  ✓ Title slide")
    
    # Slide 2: Agenda
    prs = create_agenda_slide(prs)
    print("  ✓ Agenda")
    
    # Slide 3: Algorithms
    prs = add_algorithms_slide(prs)
    print("  ✓ Algorithms")
    
    # Slides me figura
    viz_dir = 'vizualizime_doktorature'
    
    visualization_slides = [
        ("TIME SERIES ANALYSIS - Çmimet dhe Moving Averages", "01_time_series_comprehensive.png"),
        ("RETURNS DISTRIBUTION - Daily, Weekly, Monthly", "02_returns_distribution.png"),
        ("VOLATILITY ANALYSIS - Historical dhe Rolling", "03_volatility_analysis.png"),
        ("CORRELATION MATRIX - Multi-Asset Correlations", "04_correlation_heatmap.png"),
        ("PORTFOLIO PERFORMANCE - Multi-Asset Returns", "05_portfolio_performance.png"),
        ("VOLUME ANALYSIS - Trading Volumes", "06_volume_analysis.png"),
        ("CANDLESTICK CHARTS - OHLC Analysis", "07_candlestick_analysis.png"),
        ("STATISTICAL DISTRIBUTIONS - Box Plots", "08_statistical_distributions.png"),
        ("3D SURFACE PLOT - Price vs Volume vs Time", "09_3d_surface_plot.png"),
        ("RISK-RETURN SCATTER - Portfolio Optimization", "10_risk_return_scatter.png"),
        ("CUMULATIVE RETURNS - Long-term Performance", "11_cumulative_returns.png"),
        ("DRAWDOWN ANALYSIS - Risk Assessment", "12_drawdown_analysis.png"),
    ]
    
    for title, filename in visualization_slides:
        image_path = os.path.join(viz_dir, filename)
        prs = add_image_slide(prs, title, image_path)
        print(f"  ✓ {title[:50]}...")
    
    # Conclusion
    prs = create_conclusion_slide(prs)
    print("  ✓ Conclusion")
    
    # Ruaj prezantimin
    print("\n[3/4] Duke ruajtur prezantimin...")
    os.makedirs('prezantimi_powerpoint', exist_ok=True)
    output_path = 'prezantimi_powerpoint/Prezantimi_Doktoral_Analiza_Financiare.pptx'
    prs.save(output_path)
    print(f"✓ Ruajtur në: {output_path}")
    
    # Krijo VBA script për customization të mëtejshme
    print("\n[4/4] Duke krijuar VBA macro script...")
    vba_script = '''
'═══════════════════════════════════════════════════════════════════════════════
' VBA MACRO PËR CUSTOMIZATION TË POWERPOINT PREZANTIMIT
'═══════════════════════════════════════════════════════════════════════════════
' 
' Instruksione:
' 1. Hap prezantimin në PowerPoint
' 2. Shko te View > Macros
' 3. Krijo Macro të re dhe kopjo këtë kod
' 4. Ekzekuto macro-n
'
' Funksionalitetet:
' - Aplikon theme profesionale
' - Shtojnë animacione
' - Formaton tekstin
' - Shtojnë footer me numra faqesh
'═══════════════════════════════════════════════════════════════════════════════

Sub CustomizePresentation()
    Dim oSlide As Slide
    Dim oShape As Shape
    Dim slideNum As Integer
    
    ' Apliko theme
    ActivePresentation.ApplyTheme _
        "C:\\Program Files\\Microsoft Office\\root\\Document Themes 16\\Facet.thmx"
    
    ' Per çdo slide
    slideNum = 1
    For Each oSlide In ActivePresentation.Slides
        
        ' Shto footer me numer faqe (pervec slides se pare)
        If slideNum > 1 Then
            oSlide.HeadersFooters.Footer.Visible = msoTrue
            oSlide.HeadersFooters.Footer.Text = "Analiza Financiare Doktorale | " & _
                                                 "Universiteti i Prishtinës | " & _
                                                 "Faqe " & slideNum
        End If
        
        ' Apliko transition
        oSlide.SlideShowTransition.EntryEffect = ppEffectFadeSmoothly
        oSlide.SlideShowTransition.Duration = 0.5
        
        ' Animacioneeper shapes (vetem per bullet points)
        For Each oShape In oSlide.Shapes
            If oShape.HasTextFrame Then
                If oShape.TextFrame.HasText Then
                    ' Fly in animation
                    oSlide.TimeLine.MainSequence.AddEffect _
                        Shape:=oShape, _
                        effectId:=msoAnimEffectFly, _
                        trigger:=msoAnimTriggerOnPageClick
                End If
            End If
        Next oShape
        
        slideNum = slideNum + 1
    Next oSlide
    
    MsgBox "Customization e prezantimit u përfundua me sukses!", vbInformation, "Sukses"
End Sub

Sub ApplyCustomColors()
    ' Aplikon ngjyra custom per universiteti
    Dim oSlide As Slide
    Dim oShape As Shape
    
    For Each oSlide In ActivePresentation.Slides
        For Each oShape In oSlide.Shapes
            If oShape.HasTextFrame Then
                ' Ngjyra e universitetit - blu
                oShape.TextFrame.TextRange.Font.Color.RGB = RGB(0, 51, 102)
            End If
        Next oShape
    Next oSlide
    
    MsgBox "Ngjyrat u aplikuan!", vbInformation
End Sub

Sub AddCompanyLogo()
    ' Shtojnë logon e universitetit (nese ekziston)
    Dim oSlide As Slide
    Dim logoPath As String
    
    logoPath = "C:\\Path\\To\\University\\Logo.png"
    
    ' Kontrollo nese logo ekziston
    If Dir(logoPath) <> "" Then
        For Each oSlide In ActivePresentation.Slides
            oSlide.Shapes.AddPicture _
                FileName:=logoPath, _
                LinkToFile:=msoFalse, _
                SaveWithDocument:=msoTrue, _
                Left:=600, _
                Top:=10, _
                Width:=60, _
                Height:=60
        Next oSlide
        
        MsgBox "Logoja u shtua në të gjitha slides!", vbInformation
    Else
        MsgBox "Logo nuk u gjet në: " & logoPath, vbExclamation
    End If
End Sub
'''
    
    vba_path = 'prezantimi_powerpoint/CustomizePowerPoint.vba'
    with open(vba_path, 'w', encoding='utf-8') as f:
        f.write(vba_script)
    
    print(f"✓ VBA script ruajtur në: {vba_path}")
    
    # Summary
    print("\n" + "="*80)
    print("✓✓✓ POWERPOINT PREZANTIMI U GJENERUA ME SUKSES ✓✓✓")
    print("="*80)
    print(f"\nFile: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"\nPër të hapur:")
    print(f"  1. Shko te: prezantimi_powerpoint/")
    print(f"  2. Hap: Prezantimi_Doktoral_Analiza_Financiare.pptx")
    print(f"\nPër customization të mëtejshme:")
    print(f"  1. Hap prezantimin")
    print(f"  2. View > Macros")
    print(f"  3. Përdor VBA script nga: {vba_path}")
    print("="*80 + "\n")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"\n❌ GABIM: {str(e)}")
        print("\nInstalo python-pptx:")
        print("  python -m pip install python-pptx")
        import traceback
        traceback.print_exc()
